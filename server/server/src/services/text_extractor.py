import io
import tempfile
import os
import psycopg2
from core.config import settings
import pathlib
from bs4 import BeautifulSoup
# Import Google Generative AI library for fallback PDF extraction
try:
    from google import genai
except ImportError:
    genai = None

def extract_text(file_type, file_bytes_or_url):
    """
    Extracts text from file bytes or a URL based on the file type.

    Parameters:
    - file_type: A string indicating the file type (e.g., 'pdf', 'docx', 'pptx', 'img', 'url').
    - file_bytes_or_url: The file content as bytes or a URL string.

    Returns:
    - A string containing the extracted text.
    """
    file_type = file_type.lower()

    if file_type == 'pdf':
        try:
            # First attempt with PostgreSQL function
            conn = psycopg2.connect(os.getenv("DATABASE_URL"))
            with conn.cursor() as cur:
                cur.execute("SELECT rag.text_from_pdf(%s);", (psycopg2.Binary(file_bytes_or_url),))
                result = cur.fetchone()
                extracted_text = result[0] if result else ""
                print("Text extracted using Neon's rag")
            conn.close()
            
            # Check if extracted text has less than 50 words
            if len(extracted_text.split()) < 50 and genai is not None:
                print("Using Gemini for PDF extraction due to insufficient text from Neon's rag")
                # Fallback to Gemini text extractor
                try:
                    # Initialize Gemini client
                    client = genai.Client(api_key=settings.GEMINI_API_KEY)
                    
                    # Save the PDF to a temporary file
                    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
                        temp_file.write(file_bytes_or_url)
                        temp_path = pathlib.Path(temp_file.name)
                    
                    # Upload the PDF using the File API
                    sample_file = client.files.upload(file=temp_path)
                    
                    # Extract text using Gemini
                    response = client.models.generate_content(
                        model="gemini-2.0-flash",
                        contents=[sample_file, """Extract all text from this document while maintaining the original structure and flow. 

For any images, charts, or figures found in the document:
1) Provide an overall description of what the image shows
2) Extract and transcribe any text content visible within the image
3) Explain the significance of the image in context if possible

Format each image analysis as:
[IMAGE: (overall description) | TEXT IN IMAGE: (any text found in the image) | CONTEXT: (how this image relates to surrounding content)]

After extracting all content, please provide a concise summary of the entire document (2-3 sentences) under the heading [DOCUMENT SUMMARY].

Ensure that you capture the key points, main arguments, and overall purpose of the document in your extraction and summary."""]
                    )
                    
                    # Clean up the temporary file
                    temp_path.unlink()
                    
                    # Return the extracted text from Gemini
                    gemini_text = response.text
                    
                    # If Gemini also returned minimal text, use the better of the two
                    if len(gemini_text.split()) > len(extracted_text.split()):
                        return gemini_text
                    return extracted_text
                    
                except Exception as e:
                    # If Gemini extraction fails, return the original text
                    print(f"Gemini extraction failed: {str(e)}")
                    return extracted_text
            
            return extracted_text
            
        except Exception as e:
            # Fallback to pymupdf if PostgreSQL extraction fails
            print(f"PostgreSQL PDF extraction failed: {str(e)}")
            import pymupdf
            import pymupdf4llm

            # Fallback to pymupdf4llm
            bytes_io = io.BytesIO(file_bytes_or_url)
            doc = pymupdf.open(stream=bytes_io, filetype="pdf")
            md_text = pymupdf4llm.to_markdown(doc)
            return md_text
    

    elif file_type == 'docx':
        try:
            # First attempt with PostgreSQL function
            conn = psycopg2.connect(os.getenv("DATABASE_URL"))
            with conn.cursor() as cur:
                cur.execute("SELECT rag.text_from_docx(%s);", (psycopg2.Binary(file_bytes_or_url),))
                result = cur.fetchone()
                extracted_text = result[0] if result else ""
                print("Text extracted using Neon's rag")
            conn.close()
            return extracted_text
            
        except Exception as e:
            # Fallback to python-docx if PostgreSQL extraction fails
            print(f"PostgreSQL DOCX extraction failed: {str(e)}")
            import docx
            bytes_io = io.BytesIO(file_bytes_or_url)
            doc = docx.Document(bytes_io)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text

    elif file_type in ['ppt', 'pptx']:
        from pptx import Presentation
        bytes_io = io.BytesIO(file_bytes_or_url)
        prs = Presentation(bytes_io)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    elif file_type in ['img', 'jpeg', 'jpg', 'png']:
        from PIL import Image
        import pytesseract
        bytes_io = io.BytesIO(file_bytes_or_url)
        image = Image.open(bytes_io)
        text = pytesseract.image_to_string(image)
        return text

    elif file_type == 'url':
        import requests
        from urllib.parse import urlparse
        from bs4 import BeautifulSoup

        response = requests.get(file_bytes_or_url)
        response.raise_for_status()

        # Extract file extension if any
        parsed_url = urlparse(file_bytes_or_url)
        path = parsed_url.path.lower()

        if path.endswith('.pdf'):
            return extract_text('pdf', response.content)
        elif path.endswith('.docx'):
            return extract_text('docx', response.content)
        elif path.endswith(('.ppt', '.pptx')):
            return extract_text('pptx', response.content)
        elif path.endswith(('.jpg', '.jpeg', '.png')):
            return extract_text('img', response.content)

        # Fallback to HTML parsing
        soup = BeautifulSoup(response.content, 'html.parser')
        for tag in soup(['script', 'style']):
            tag.decompose()
        
        text = soup.get_text(separator='\n')
        return '\n'.join(line.strip() for line in text.splitlines() if line.strip())
    
    elif file_type == 'json':
        import json

        if isinstance(file_bytes_or_url, (bytes, bytearray)):
            data = json.loads(file_bytes_or_url.decode('utf-8'))
        elif isinstance(file_bytes_or_url, str):
            # Assume it's a JSON string (from a URL maybe)
            data = json.loads(file_bytes_or_url)
        else:
            raise ValueError("Unsupported JSON input type")

        def flatten_json(obj, prefix=''):
            lines = []
            if isinstance(obj, dict):
                for k, v in obj.items():
                    lines.extend(flatten_json(v, f"{prefix}{k}: "))
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    lines.extend(flatten_json(item, f"{prefix}[{i}] "))
            else:
                lines.append(f"{prefix}{str(obj)}")
            return lines

        flat_lines = flatten_json(data)
        return "\n".join(flat_lines)

    elif file_type == 'csv':
        import csv

        if isinstance(file_bytes_or_url, bytes):
            csv_text = file_bytes_or_url.decode('utf-8')
        else:
            csv_text = file_bytes_or_url  # maybe string from URL

        reader = csv.DictReader(csv_text.splitlines())
        rows = []

        for i, row in enumerate(reader, 1):
            lines = [f"{k.strip()}: {v.strip() if v.strip() else 'No Data'}" for k, v in row.items()]
            rows.append(f"Row {i}:\n" + "\n".join(lines))

        return "\n\n".join(rows)



    else:
        raise ValueError("Unsupported file type")
